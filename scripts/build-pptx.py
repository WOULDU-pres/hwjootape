#!/usr/bin/env python3
"""build-pptx.py — assemble an editable .pptx from a resolved deck spec.

Usage: build-pptx.py <spec.json> <output.pptx>

The spec is RESOLVED (asset paths absolute, text final). Coordinates are
normalized [0,1] relative to the slide; this script maps them to EMU.
Text elements become real, editable PowerPoint text boxes (selectable text);
image elements become pictures. Korean fonts are set for both the Latin and
East-Asian typeface slots so Hangul renders correctly.

Resolved spec shape:
{
  "slideWidthEmu": 12192000, "slideHeightEmu": 6858000,   # optional, defaults to 16:9
  "slides": [
    {
      "background": { "path": "/abs/bg.png" } | null,
      "elements": [
        { "type":"text", "text":"제목", "nbbox":{"x":..,"y":..,"w":..,"h":..},
          "color":"1a3a8f", "fontSizePt":40, "bold":true, "align":"left",
          "font":"Pretendard", "fontEA":"Apple SD Gothic Neo" },
        { "type":"image", "path":"/abs/el.png", "nbbox":{...}, "z":2 }
      ]
    }
  ]
}
A bare single-slide spec (no "slides" key, has "elements") is also accepted.
"""
import json
import os
import sys

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

DEFAULT_W_EMU = 12192000  # 13.333"
DEFAULT_H_EMU = 6858000   # 7.5"  -> 16:9
EMU_PER_PT = 12700

ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

# Bullet marker compose emits per logical line (see src/lib/slides/compose.ts):
# "• " (U+2022 + space). A bullet paragraph gets a real hanging indent so that
# when PowerPoint wraps a long bullet, the continuation lines align under the
# bullet BODY instead of dropping to the box's left edge.
BULLET_PREFIX = "• "

_KOREAN_TTF = [
    ("/System/Library/Fonts/AppleSDGothicNeo.ttc", 0),
    ("/System/Library/Fonts/Supplemental/AppleGothic.ttf", 0),
]


def _fit_font_pt(text, start_pt, box_w_pt, box_h_pt, min_pt=8):
    """Largest point size at which `text` wraps within the box (width AND height).
    Uses Pillow to measure so re-rendered Korean never overflows the OCR-derived
    bbox (R5). Falls back to start_pt if Pillow/fonts are unavailable."""
    try:
        from PIL import ImageFont
    except Exception:
        return start_pt
    fpath, fidx = next(((p, i) for p, i in _KOREAN_TTF if os.path.exists(p)), (None, 0))
    if not fpath:
        return start_pt

    def wrap(words_font, max_w):
        # Mirror render-png.py: bullet continuation lines hang-indent, so they
        # have less width available; account for that when measuring fit.
        lines = []
        for raw in str(text).split("\n"):
            indent_w = words_font.getlength(BULLET_PREFIX) if raw.startswith(BULLET_PREFIX) else 0.0
            cur, started = "", False
            for w in raw.split(" "):
                avail = max_w - (indent_w if started else 0)
                trial = w if not cur else cur + " " + w
                if words_font.getlength(trial) <= avail or not cur:
                    cur = trial
                else:
                    lines.append(cur); cur = w; started = True
            lines.append(cur); started = True
        return lines or [""]

    pt = max(min_pt, int(start_pt))
    while pt >= min_pt:
        try:
            f = ImageFont.truetype(fpath, pt, index=fidx)
        except Exception:
            return start_pt
        lines = wrap(f, box_w_pt)
        widest = max((f.getlength(ln) for ln in lines), default=0)
        line_h = (f.getbbox("Ag")[3] - f.getbbox("Ag")[1]) * 1.25
        if widest <= box_w_pt and line_h * len(lines) <= box_h_pt:
            return pt
        pt = int(pt * 0.92) if pt > min_pt else min_pt - 1
    return min_pt


def die(msg: str) -> "NoReturn":
    sys.stderr.write(f"build-pptx: {msg}\n")
    sys.exit(2)


def emu_box(nbbox, w_emu, h_emu):
    return (
        Emu(int(round(nbbox["x"] * w_emu))),
        Emu(int(round(nbbox["y"] * h_emu))),
        Emu(int(round(nbbox["w"] * w_emu))),
        Emu(int(round(nbbox["h"] * h_emu))),
    )


def _marker_emu(size_pt):
    """Width of the '• ' bullet marker, in EMU, at the given point size — used as
    the paragraph hanging-indent amount. Falls back to a geometric estimate when
    Pillow/fonts are unavailable."""
    try:
        from PIL import ImageFont
        fpath, fidx = next(((p, i) for p, i in _KOREAN_TTF if os.path.exists(p)), (None, 0))
        if fpath:
            f = ImageFont.truetype(fpath, max(1, int(round(size_pt))), index=fidx)
            return int(round(f.getlength(BULLET_PREFIX) * EMU_PER_PT))
    except Exception:
        pass
    # '•' ~0.5em + space ~0.27em -> ~0.77em wide
    return int(round(size_pt * 0.77 * EMU_PER_PT))


def set_hanging_indent(p, marker_emu):
    """Give paragraph `p` a real hanging indent: the first line starts at the box
    edge (where '• ' sits) and wrapped lines hang at the marker width, aligning
    under the bullet body. marL = markerWidth, indent = -markerWidth."""
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(max(0, int(marker_emu))))
    pPr.set("indent", str(-max(0, int(marker_emu))))


def set_ea_font(run, typeface: str):
    """python-pptx exposes only the Latin font; set the East-Asian slot via XML."""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        existing = rPr.find(qn(tag))
        if existing is None:
            el = rPr.makeelement(qn(tag), {"typeface": typeface})
            rPr.append(el)
        else:
            existing.set("typeface", typeface)


def add_text(slide, el, w_emu, h_emu):
    left, top, width, height = emu_box(el["nbbox"], w_emu, h_emu)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    # zero internal margins so the box matches the OCR bbox
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    # PowerPoint shrinks text to the shape; we also pre-fit the size below so
    # static viewers (Keynote/LibreOffice) that ignore autofit still look right.
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass

    lines = str(el.get("text", "")).split("\n")
    font_name = el.get("font", "Pretendard")
    font_ea = el.get("fontEA", el.get("font", "Apple SD Gothic Neo"))
    requested = el.get("fontSizePt", 24)
    size = _fit_font_pt(el.get("text", ""), requested, width / EMU_PER_PT, height / EMU_PER_PT)
    bold = bool(el.get("bold", False))
    color = el.get("color")
    align = ALIGN.get(el.get("align", "left"), PP_ALIGN.LEFT)

    marker_emu = _marker_emu(size)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        # Bullet line: real hanging indent so wrapped continuation lines align
        # under the bullet body instead of dropping to the box's left edge.
        if line.startswith(BULLET_PREFIX):
            set_hanging_indent(p, marker_emu)
        run = p.add_run()
        run.text = line
        f = run.font
        f.size = Pt(size)
        f.bold = bold
        f.name = font_name
        set_ea_font(run, font_ea)
        if color:
            try:
                f.color.rgb = RGBColor.from_string(str(color).lstrip("#"))
            except Exception:
                pass


def _parse_hex(c):
    """Hex (leading '#' tolerated) -> RGBColor, or None when absent/invalid."""
    if not c:
        return None
    try:
        return RGBColor.from_string(str(c).lstrip("#"))
    except Exception:
        return None


def add_solid_fill(slide, rgb, w_emu, h_emu):
    """Paint a full-bleed borderless rectangle behind everything as the slide's
    solid background, so the .pptx canvas color matches the PNG. (Editing the
    slide-master background is brittle across viewers; a full-slide shape is the
    robust, WYSIWYG equivalent.)"""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(w_emu), Emu(h_emu))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb
    shp.line.fill.background()
    shp.shadow.inherit = False


def add_accent_bar(slide, bar, w_emu, h_emu):
    """Paint the theme accent band as a borderless filled rectangle at its nbbox."""
    rgb = _parse_hex(bar.get("color"))
    if rgb is None or not bar.get("nbbox"):
        return
    left, top, width, height = emu_box(bar["nbbox"], w_emu, h_emu)
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb
    shp.line.fill.background()
    shp.shadow.inherit = False


def add_image(slide, el, w_emu, h_emu):
    left, top, width, height = emu_box(el["nbbox"], w_emu, h_emu)
    slide.shapes.add_picture(el["path"], left, top, width=width, height=height)


def main():
    if len(sys.argv) < 3:
        die("usage: build-pptx.py <spec.json> <output.pptx>")
    spec_path, out_path = sys.argv[1], sys.argv[2]
    with open(spec_path, encoding="utf-8") as fh:
        spec = json.load(fh)

    w_emu = int(spec.get("slideWidthEmu", DEFAULT_W_EMU))
    h_emu = int(spec.get("slideHeightEmu", DEFAULT_H_EMU))
    slides = spec.get("slides")
    if slides is None:
        slides = [spec] if "elements" in spec else []
    if not slides:
        die("spec has no slides")

    prs = Presentation()
    prs.slide_width = Emu(w_emu)
    prs.slide_height = Emu(h_emu)
    blank = prs.slide_layouts[6]  # blank

    for s in slides:
        slide = prs.slides.add_slide(blank)
        # Solid theme fill behind everything (absent => leave the default white
        # master background untouched, the historical behavior).
        fill_rgb = _parse_hex(s.get("backgroundColor"))
        if fill_rgb is not None:
            add_solid_fill(slide, fill_rgb, w_emu, h_emu)
        bg = s.get("background")
        if bg and bg.get("path"):
            # A full-slide background PICTURE wins over the solid fill.
            slide.shapes.add_picture(bg["path"], Emu(0), Emu(0), width=Emu(w_emu), height=Emu(h_emu))
        # Theme accent band above the fill/picture, behind the elements.
        bar = s.get("accentBar")
        if bar:
            add_accent_bar(slide, bar, w_emu, h_emu)
        for el in sorted(s.get("elements", []), key=lambda e: e.get("z", 0)):
            t = el.get("type")
            if t == "text":
                add_text(slide, el, w_emu, h_emu)
            elif t == "image":
                add_image(slide, el, w_emu, h_emu)

    prs.save(out_path)
    sys.stderr.write(f"build-pptx: wrote {len(slides)} slide(s) -> {out_path}\n")


if __name__ == "__main__":
    main()
